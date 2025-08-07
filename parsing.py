"""
Document parsing system for PDF, PPTX, and Excel files.
"""
import re
import os
import pandas as pd
from typing import List, Dict, Any
from pathlib import Path
from pptx import Presentation
from pypdf import PdfReader
import pdfplumber


class PDFParser:
    """PDF text extraction with number formatting correction."""
    
    def parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text and tables from PDF file."""
        try:
            filename = Path(file_path).name
            pages = []
            
            try:
                pages = self._extract_with_pdfplumber(file_path)
                print(f"Successfully parsed PDF with pdfplumber: {filename}")
            except Exception as e:
                print(f"pdfplumber failed, trying pypdf: {e}")
                try:
                    pages = self._extract_with_pypdf(file_path)
                    print(f"Successfully parsed PDF with pypdf: {filename}")
                except Exception as e2:
                    print(f"Both PDF extraction methods failed: {e2}")
                    return self._create_fallback_report(file_path)
            
            processed_pages = []
            for page_data in pages:
                processed_text = self._post_process_text(page_data['text'])
                if processed_text.strip():
                    processed_pages.append({
                        'page': page_data['page'],
                        'text': processed_text.strip()
                    })
            
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(processed_pages),
                    'text_blocks_amount': len(processed_pages),
                    'tables_amount': 0,
                    'pictures_amount': 0,
                    'document_type': 'pdf'
                },
                'content': {'pages': processed_pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed PDF: {filename} ({len(processed_pages)} pages)")
            return report
            
        except Exception as e:
            print(f"Error parsing PDF file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _extract_with_pdfplumber(self, file_path: str) -> List[Dict]:
        """Extract text and tables using pdfplumber with layout detection."""
        pages = []
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text_parts = []
                
                # Perform layout analysis on the page
                layout_info = self._analyze_page_layout(page)
                
                # Extract text based on layout type
                if layout_info['layout_type'] == 'multi_column':
                    # Multi-column extraction with coordinate-based column detection
                    column_texts = self._extract_multi_column_text(page, layout_info)
                    if column_texts:
                        text_parts.extend(column_texts)
                else:
                    # Standard single-column extraction
                    standard_text = page.extract_text()
                    if standard_text:
                        text_parts.append(standard_text)
                
                # Extract layout-preserved text for complex layouts
                try:
                    layout_text = page.extract_text(layout=True, x_tolerance=1, y_tolerance=1)
                    if layout_text and layout_text not in text_parts:
                        text_parts.append("=== Layout Preserved ===")
                        text_parts.append(layout_text)
                except:
                    pass
                
                # Extract tables
                try:
                    tables = page.extract_tables()
                    if tables:
                        text_parts.append("=== Tables ===")
                        for i, table in enumerate(tables):
                            if table:
                                table_text = self._format_table_text(table)
                                text_parts.append(f"Table {i+1}:\n{table_text}")
                except:
                    pass
                
                combined_text = '\n\n'.join(text_parts)
                if combined_text.strip():
                    page_data = {
                        'page': page_num,
                        'text': combined_text.strip(),
                        'layout_type': layout_info['layout_type'],
                        'column_count': layout_info['column_count'],
                        'extraction_method': layout_info['extraction_method']
                    }
                    pages.append(page_data)
        
        return pages
    
    def _analyze_page_layout(self, page) -> Dict[str, Any]:
        """Analyze PDF page layout to detect columns and structure."""
        try:
            # Get page dimensions
            page_width = page.width
            page_height = page.height
            
            # Extract characters with coordinates for layout analysis
            chars = page.chars
            if not chars:
                return {
                    'layout_type': 'single_column',
                    'column_count': 1,
                    'extraction_method': 'standard'
                }
            
            # Analyze character distribution across page width
            x_positions = [char['x0'] for char in chars if 'x0' in char]
            if not x_positions:
                return {
                    'layout_type': 'single_column',
                    'column_count': 1,
                    'extraction_method': 'standard'
                }
            
            # Group characters by approximate x-position to detect columns
            x_bins = self._create_position_bins(x_positions, page_width)
            column_gaps = self._detect_column_gaps(x_bins, page_width)
            
            # Determine layout type based on gaps and distribution
            if len(column_gaps) >= 1:
                # Check if gaps are significant enough to indicate columns
                significant_gaps = [gap for gap in column_gaps if gap['width'] > page_width * 0.02]
                
                if significant_gaps:
                    column_count = len(significant_gaps) + 1
                    layout_type = 'multi_column' if column_count > 1 else 'single_column'
                else:
                    column_count = 1
                    layout_type = 'single_column'
            else:
                column_count = 1
                layout_type = 'single_column'
            
            # Additional complexity check based on text object distribution
            text_objects = page.objects.get('char', [])
            if len(text_objects) > 1000 and layout_type == 'single_column':
                layout_type = 'complex'
            
            return {
                'layout_type': layout_type,
                'column_count': column_count,
                'extraction_method': 'coordinate_based',
                'page_width': page_width,
                'column_gaps': column_gaps if 'significant_gaps' in locals() else []
            }
            
        except Exception as e:
            print(f"Warning: Layout analysis failed for page, using default: {e}")
            return {
                'layout_type': 'single_column',
                'column_count': 1,
                'extraction_method': 'standard'
            }
    
    def _create_position_bins(self, x_positions: List[float], page_width: float, bin_count: int = 50) -> List[int]:
        """Create bins for x-position distribution analysis."""
        if not x_positions:
            return []
        
        # Create histogram bins across page width
        bin_width = page_width / bin_count
        bins = [0] * bin_count
        
        for x in x_positions:
            bin_index = min(int(x / bin_width), bin_count - 1)
            bins[bin_index] += 1
        
        return bins
    
    def _detect_column_gaps(self, x_bins: List[int], page_width: float) -> List[Dict]:
        """Detect gaps between columns based on character distribution."""
        if not x_bins:
            return []
        
        # Find regions with very low character density (potential gaps)
        total_chars = sum(x_bins)
        if total_chars == 0:
            return []
        
        # Calculate moving average to smooth the distribution
        window_size = 3
        smoothed_bins = []
        for i in range(len(x_bins)):
            start_idx = max(0, i - window_size // 2)
            end_idx = min(len(x_bins), i + window_size // 2 + 1)
            avg = sum(x_bins[start_idx:end_idx]) / (end_idx - start_idx)
            smoothed_bins.append(avg)
        
        # Identify gaps (regions with density below threshold)
        avg_density = total_chars / len(x_bins)
        gap_threshold = avg_density * 0.1  # 10% of average density
        
        gaps = []
        in_gap = False
        gap_start = 0
        bin_width = page_width / len(x_bins)
        
        for i, density in enumerate(smoothed_bins):
            if density <= gap_threshold and not in_gap:
                # Start of gap
                in_gap = True
                gap_start = i
            elif density > gap_threshold and in_gap:
                # End of gap
                in_gap = False
                gap_width = (i - gap_start) * bin_width
                if gap_width > page_width * 0.02:  # Minimum 2% of page width
                    gaps.append({
                        'start': gap_start * bin_width,
                        'end': i * bin_width,
                        'width': gap_width
                    })
        
        # Handle gap that extends to end of page
        if in_gap:
            gap_width = (len(smoothed_bins) - gap_start) * bin_width
            if gap_width > page_width * 0.02:
                gaps.append({
                    'start': gap_start * bin_width,
                    'end': page_width,
                    'width': gap_width
                })
        
        return gaps
    
    def _extract_multi_column_text(self, page, layout_info: Dict) -> List[str]:
        """Extract text from multi-column layout using coordinate-based approach."""
        try:
            column_gaps = layout_info.get('column_gaps', [])
            page_width = layout_info.get('page_width', page.width)
            
            if not column_gaps:
                # Fallback to standard extraction
                standard_text = page.extract_text()
                return [standard_text] if standard_text else []
            
            # Define column boundaries based on detected gaps
            column_boundaries = [0]  # Start with left edge
            for gap in column_gaps:
                column_boundaries.append(gap['start'])
                column_boundaries.append(gap['end'])
            column_boundaries.append(page_width)  # End with right edge
            
            # Remove duplicates and sort
            column_boundaries = sorted(list(set(column_boundaries)))
            
            # Create column regions (pairs of boundaries)
            columns = []
            for i in range(0, len(column_boundaries) - 1, 2):
                if i + 1 < len(column_boundaries):
                    left = column_boundaries[i]
                    right = column_boundaries[i + 1]
                    columns.append((left, right))
            
            # Extract text from each column
            column_texts = []
            for i, (left, right) in enumerate(columns):
                # Create a cropped version of the page for this column
                cropped_page = page.within_bbox((left, 0, right, page.height))
                column_text = cropped_page.extract_text()
                
                if column_text and column_text.strip():
                    column_texts.append(f"||COLUMN|| {column_text}")
            
            return column_texts if column_texts else [page.extract_text()]
            
        except Exception as e:
            print(f"Warning: Multi-column extraction failed, using standard: {e}")
            standard_text = page.extract_text()
            return [standard_text] if standard_text else []
    
    def _extract_with_pypdf(self, file_path: str) -> List[Dict]:
        """Extract text using pypdf as fallback."""
        pages = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append({'page': page_num, 'text': text.strip()})
        
        return pages
    
    def _format_table_text(self, table: List[List]) -> str:
        """Format extracted table data into readable text."""
        if not table:
            return ""
        
        formatted_rows = []
        for row in table:
            if row:
                clean_cells = []
                for cell in row:
                    cell_text = str(cell).strip() if cell is not None else ""
                    clean_cells.append(cell_text)
                formatted_rows.append(" | ".join(clean_cells))
        
        return "\n".join(formatted_rows)
    
    def _post_process_text(self, text: str) -> str:
        """Fix number formatting and spacing issues in extracted text."""
        if not text:
            return text
        
        processed_text = self._fix_number_formatting(text)
        processed_text = self._fix_spacing_issues(processed_text)
        processed_text = self._fix_symbols(processed_text)
        
        return processed_text
    
    def _fix_number_formatting(self, text: str) -> str:
        """Correct decimal points, thousand separators, and currency symbols."""
        text = re.sub(r'(\d)\s*[,，]\s*(\d{3})', r'\1,\2', text)
        text = re.sub(r'(\d)\s*[.．]\s*(\d)', r'\1.\2', text)
        text = re.sub(r'(\d)\s*[oO]\s*(\d)', r'\1.0\2', text)
        text = re.sub(r'(\d)\s+(\d)', r'\1\2', text)
        text = re.sub(r'(\d)\s*[%％]', r'\1%', text)
        text = re.sub(r'[$＄]\s*(\d)', r'$\1', text)
        text = re.sub(r'([NT$]+)\s*(\d)', r'\1\2', text)
        text = re.sub(r'[-－—]\s*(\d)', r'-\1', text)
        
        return text
    
    def _fix_spacing_issues(self, text: str) -> str:
        """Fix spacing and line break issues."""
        text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
        text = re.sub(r' {2,}', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        return text
    
    def _fix_symbols(self, text: str) -> str:
        """Fix currency and other symbols."""
        text = re.sub(r'[＄$]', '$', text)
        text = re.sub(r'[％%]', '%', text)
        text = re.sub(r'[（(]', '(', text)
        text = re.sub(r'[）)]', ')', text)
        
        return text
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when PDF parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }


class PPTXParser:
    """PPTX content extraction including tables, charts, and images."""
    
    def parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX file for all types of content."""
        try:
            filename = Path(file_path).name
            prs = Presentation(file_path)
            
            pages = []
            tables_found = 0
            charts_found = 0
            images_found = 0
            other_objects_found = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = self._extract_slide_content(slide)
                combined_text = slide_content['combined_text']
                stats = slide_content['stats']
                
                # Update counts
                tables_found += stats['tables']
                charts_found += stats['charts']
                images_found += stats['images']
                other_objects_found += stats['other_objects']
                
                if combined_text.strip():
                    pages.append({
                        'page': slide_num,
                        'text': combined_text.strip()
                    })
            
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(pages),
                    'text_blocks_amount': len(pages),
                    'tables_amount': tables_found,
                    'pictures_amount': images_found,
                    'document_type': 'pptx'
                },
                'content': {'pages': pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed PPTX: {filename}")
            print(f"  - {len(pages)} slides with content")
            print(f"  - {tables_found} tables, {charts_found} charts, {images_found} images")
            
            return report
            
        except Exception as e:
            print(f"Error parsing PPTX file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _extract_slide_content(self, slide) -> Dict[str, Any]:
        """Extract all types of content from a single slide."""
        content_parts = []
        stats = {'tables': 0, 'charts': 0, 'images': 0, 'other_objects': 0}
        
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                shape_content = self._process_shape(shape, stats)
                if shape_content:
                    content_parts.append(shape_content)
        
        return {
            'combined_text': '\n\n'.join(content_parts),
            'stats': stats
        }
    
    def _process_shape(self, shape, stats: Dict[str, int]) -> str:
        """Process individual shape and extract relevant content."""
        content_parts = []
        
        try:
           
            shape_type_info = ""
            if hasattr(shape, 'shape_type'):
                shape_type_info = f" (type: {shape.shape_type})"
            
            if self._has_table(shape):
                stats['tables'] += 1
                table_content = self._extract_table_text(shape.table)
                content_parts.append(f"Table{shape_type_info}:\n{table_content}")
            
            elif self._has_chart(shape):
                stats['charts'] += 1
                chart_content = self._extract_chart_text(shape.chart)
                content_parts.append(f"Chart{shape_type_info}:\n{chart_content}")
            
            elif self._is_image_shape(shape):
                stats['images'] += 1
                image_info = self._extract_image_info(shape)
                content_parts.append(f"Image{shape_type_info}:\n{image_info}")
            
            elif self._is_group_shape(shape):
                group_content = self._extract_group_content(shape, stats)
                if group_content:
                    content_parts.append(f"Group{shape_type_info}:\n{group_content}")
            
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                text_content = self._extract_text_frame(shape.text_frame)
                if text_content:
                    content_parts.append(f"Text Frame{shape_type_info}:\n{text_content}")
            
            else:
                stats['other_objects'] += 1
                other_content = self._extract_other_shape_content(shape)
                if other_content:
                    content_parts.append(f"Other{shape_type_info}:\n{other_content}")
        
        except Exception as e:
            print(f"Warning: Error processing shape: {e}")
            stats['other_objects'] += 1
        
        return '\n\n'.join(content_parts)
    
    def _has_table(self, shape) -> bool:
        try:
            # python-pptx throws ValueError if shape doesn't contain a table
            table = shape.table
            return table is not None
        except (ValueError, AttributeError):
            return False
    
    def _has_chart(self, shape) -> bool:
        try:
            # python-pptx throws ValueError if shape doesn't contain a chart
            chart = shape.chart
            return chart is not None
        except (ValueError, AttributeError):
            return False
    
    def _is_image_shape(self, shape) -> bool:
        try:
            return hasattr(shape, 'image') or 'Picture' in str(type(shape))
        except:
            return False
    
    def _is_group_shape(self, shape) -> bool:
        try:
            return hasattr(shape, 'shapes') and shape.shapes is not None
        except:
            return False
    
    def _extract_table_text(self, table) -> str:
        """Extract text from PPTX table."""
        if not table:
            return ""
        
        table_text = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                row_text.append(cell_text)
            table_text.append(" | ".join(row_text))
        
        return "\n".join(table_text)
    
    def _extract_chart_text(self, chart) -> str:
        """Extract basic information from PPTX chart."""
        try:
            chart_info = []
            
            # Try to get chart title
            if hasattr(chart, 'chart_title') and chart.chart_title:
                try:
                    if hasattr(chart.chart_title, 'text_frame') and chart.chart_title.text_frame:
                        title_text = chart.chart_title.text_frame.text
                        if title_text:
                            chart_info.append(f"Title: {title_text}")
                except:
                    pass
            
            # Try to get chart type
            if hasattr(chart, 'chart_type'):
                chart_info.append(f"Chart Type: {str(chart.chart_type)}")
            
            # Try to extract some data if available
            if hasattr(chart, 'series') and chart.series:
                chart_info.append(f"Data Series: {len(chart.series)} series")
                
                # Try to extract series names and some data points
                for i, series in enumerate(chart.series):
                    try:
                        if hasattr(series, 'name') and series.name:
                            chart_info.append(f"Series {i+1}: {series.name}")
                    except:
                        pass
            
            # Try to get category and value axes information
            if hasattr(chart, 'category_axis') and chart.category_axis:
                try:
                    if hasattr(chart.category_axis, 'axis_title') and chart.category_axis.axis_title:
                        cat_title = chart.category_axis.axis_title.text_frame.text
                        if cat_title:
                            chart_info.append(f"Category Axis: {cat_title}")
                except:
                    pass
            
            if hasattr(chart, 'value_axis') and chart.value_axis:
                try:
                    if hasattr(chart.value_axis, 'axis_title') and chart.value_axis.axis_title:
                        val_title = chart.value_axis.axis_title.text_frame.text
                        if val_title:
                            chart_info.append(f"Value Axis: {val_title}")
                except:
                    pass
            
            return "\n".join(chart_info) if chart_info else "Chart (basic info not extractable)"
            
        except Exception as e:
            return f"Chart (extraction error: {str(e)[:50]}...)"
    
    def _extract_image_info(self, shape) -> str:
        """Extract basic information about image."""
        try:
            info = []
            if hasattr(shape, 'name') and shape.name:
                info.append(f"Name: {shape.name}")
            
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                info.append(f"Dimensions: {shape.width} x {shape.height}")
            
            return "\n".join(info) if info else "Image"
        except:
            return "Image"
    
    def _extract_group_content(self, group_shape, stats: Dict[str, int]) -> str:
        """Extract content from grouped shapes."""
        content_parts = []
        
        try:
            if hasattr(group_shape, 'shapes'):
                for shape in group_shape.shapes:
                    shape_content = self._process_shape(shape, stats)
                    if shape_content:
                        content_parts.append(shape_content)
        except:
            pass
        
        return '\n\n'.join(content_parts)
    
    def _extract_text_frame(self, text_frame) -> str:
        """Extract text from text frame."""
        try:
            return text_frame.text.strip() if text_frame.text else ""
        except:
            return ""
    
    def _extract_other_shape_content(self, shape) -> str:
        """Extract content from other types of shapes."""
        try:
            if hasattr(shape, 'text') and shape.text:
                return shape.text.strip()
            elif hasattr(shape, 'name') and shape.name:
                return f"Object: {shape.name}"
        except:
            pass
        
        return ""
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when PPTX parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }


class ExcelParser:
    """Excel file parsing for all sheets and data."""
    
    def parse_excel(self, file_path: str) -> Dict[str, Any]:
        """Parse Excel file for all sheets."""
        try:
            filename = Path(file_path).name
            
            # Read all sheets
            sheets_dict = pd.read_excel(file_path, sheet_name=None, header=None)
            
            pages = []
            total_rows = 0
            
            for sheet_name, df in sheets_dict.items():
                if df.empty:
                    continue
                
                # Convert DataFrame to text representation
                sheet_text_parts = [f"Sheet: {sheet_name}"]
                
                # Convert to string with proper formatting
                df_str = df.to_string(index=False, header=False, na_rep='')
                sheet_text_parts.append(df_str)
                
                combined_sheet_text = '\n'.join(sheet_text_parts)
                
                if combined_sheet_text.strip():
                    pages.append({
                        'page': len(pages) + 1,
                        'text': combined_sheet_text.strip()
                    })
                    total_rows += len(df)
            
            report = {
                'metainfo': {
                    'sha1_name': filename.rsplit('.', 1)[0],
                    'filename': filename,
                    'pages_amount': len(pages),
                    'text_blocks_amount': len(pages),
                    'tables_amount': len(sheets_dict),
                    'pictures_amount': 0,
                    'document_type': 'excel'
                },
                'content': {'pages': pages},
                'tables': [],
                'pictures': []
            }
            
            print(f"Successfully parsed Excel: {filename}")
            print(f"  - {len(pages)} sheets with {total_rows} total rows")
            
            return report
            
        except Exception as e:
            print(f"Error parsing Excel file {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when Excel parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }


class UnifiedDocumentParser:
    """Route documents to appropriate parser by file extension."""
    
    def __init__(self):
        self.pdf_parser = PDFParser()
        self.pptx_parser = PPTXParser()
        self.excel_parser = ExcelParser()
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """Parse document based on file extension."""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext == '.pdf':
                return self.pdf_parser.parse_pdf(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self.pptx_parser.parse_pptx(file_path)
            elif file_ext in ['.xls', '.xlsx']:
                return self.excel_parser.parse_excel(file_path)
            else:
                print(f"Unsupported file type: {file_ext}")
                return self._create_fallback_report(file_path)
                
        except Exception as e:
            print(f"Error parsing document {file_path}: {e}")
            return self._create_fallback_report(file_path)
    
    def _create_fallback_report(self, file_path: str) -> Dict[str, Any]:
        """Create minimal report when parsing fails."""
        filename = Path(file_path).name
        return {
            'metainfo': {
                'sha1_name': filename.rsplit('.', 1)[0],
                'filename': filename,
                'pages_amount': 0,
                'text_blocks_amount': 0,
                'tables_amount': 0,
                'pictures_amount': 0,
                'document_type': 'failed'
            },
            'content': {'pages': []},
            'tables': [],
            'pictures': []
        }
