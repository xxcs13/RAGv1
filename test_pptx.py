#!/usr/bin/env python3
"""
簡單測試PPTX解析器
"""

def test_basic_pptx():
    from pptx import Presentation
    prs = Presentation('/home/xcs/intern/rag/0.pptx')
    print(f"Found {len(prs.slides)} slides")
    return True

def test_shape_detection():
    from parsing import PPTXParser
    
    class SimplePPTXParser(PPTXParser):
        def _has_table(self, shape):
            try:
                table = shape.table
                return table is not None
            except:
                return False
        
        def _has_chart(self, shape):
            try:
                chart = shape.chart
                return chart is not None  
            except:
                return False
    
    parser = SimplePPTXParser()
    
    # Test just one slide
    from pptx import Presentation
    prs = Presentation('/home/xcs/intern/rag/0.pptx')
    slide = prs.slides[1]  # Second slide which has a table
    
    tables = 0
    charts = 0
    
    for shape in slide.shapes:
        if parser._has_table(shape):
            tables += 1
            print(f"Found table!")
        if parser._has_chart(shape):
            charts += 1
            print(f"Found chart!")
    
    print(f"Slide 2: {tables} tables, {charts} charts")

if __name__ == "__main__":
    print("Testing basic PPTX loading...")
    test_basic_pptx()
    
    print("Testing shape detection...")
    test_shape_detection()
